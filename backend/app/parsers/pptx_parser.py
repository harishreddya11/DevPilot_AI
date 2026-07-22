from pptx import Presentation

from app.parsers.base_parser import BaseParser


class PowerPointParser(BaseParser):
    """
    Parser for PowerPoint (.pptx) files.
    """

    def extract_text(
        self,
        file_path: str,
    ) -> str:

        presentation = Presentation(file_path)

        slides = []

        for index, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            slides.append(f"Slide {index}")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        slides.append(text)

        return "\n".join(slides)